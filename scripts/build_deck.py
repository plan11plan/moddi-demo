"""모디 기획서 구체화 PPT 빌더 — 웜 브랜드 테마. 실행: .venv/bin/python scripts/build_deck.py"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE

# ---- palette ----
ESPRESSO = RGBColor(0x2A, 0x24, 0x22)
ESPRESSO_SOFT = RGBColor(0x3A, 0x32, 0x2E)
CREAM = RGBColor(0xF5, 0xEF, 0xE6)
CREAM_SOFT = RGBColor(0xFB, 0xF7, 0xF0)
TERRA = RGBColor(0xE0, 0x67, 0x3C)
TERRA_DK = RGBColor(0xC5, 0x55, 0x2E)
MUTED = RGBColor(0x8B, 0x80, 0x79)
LINE = RGBColor(0xE6, 0xDD, 0xD1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_DIM = RGBColor(0xC9, 0xBE, 0xAE)
FONT = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
MX = Inches(0.92)  # left margin
CW = SW - MX * 2   # content width

_page = 0

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.06, space_after=2):
    """runs: list of (str, size_pt, color, bold) or list-of-lines (each a list of runs)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (s, size, color, bold) in line:
            r = p.add_run(); r.text = s
            r.font.size = Pt(size); r.font.name = FONT
            r.font.bold = bold; r.font.color.rgb = color
    return tb

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=1.0, radius=0.09,
         shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if radius is not None:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp

def chip(slide, l, t, w, label, fill=TERRA, fg=CREAM, size=12, h=Inches(0.34)):
    sp = rect(slide, l, t, w, h, fill=fill, radius=0.5)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(size); r.font.name = FONT
    r.font.bold = True; r.font.color.rgb = fg
    return sp

def base(title, kicker=None, dark=False):
    """표준 콘텐츠 슬라이드: kicker + 타이틀 + 푸터."""
    global _page
    _page += 1
    s = prs.slides.add_slide(BLANK)
    bg(s, ESPRESSO if dark else CREAM)
    title_col = CREAM if dark else ESPRESSO
    if kicker:
        text(s, MX, Inches(0.5), CW, Inches(0.34),
             [(kicker, 13, TERRA if not dark else TERRA, True)])
    text(s, MX, Inches(0.82), CW, Inches(0.9), [(title, 30, title_col, True)])
    # accent underline
    rect(s, MX, Inches(1.55), Inches(0.7), Pt(4), fill=TERRA, radius=0.5)
    # footer
    text(s, MX, SH - Inches(0.46), Inches(5), Inches(0.3),
         [[("모디", 9, MUTED if not dark else CREAM_DIM, True),
           (".", 9, TERRA, True),
           ("  ·  기획서 구체화", 9, MUTED if not dark else CREAM_DIM, False)]])
    text(s, SW - MX - Inches(1), SH - Inches(0.46), Inches(1), Inches(0.3),
         [(str(_page), 9, MUTED if not dark else CREAM_DIM, False)], align=PP_ALIGN.RIGHT)
    return s

def body_top():
    return Inches(1.85)

# ============================================================ 1. TITLE
s = prs.slides.add_slide(BLANK); bg(s, ESPRESSO)
rect(s, Inches(0), Inches(0), Inches(0.32), SH, fill=TERRA, radius=0)
text(s, MX, Inches(2.2), CW, Inches(1.2),
     [[("모디", 60, CREAM, True), (".", 60, TERRA, True)]])
text(s, MX, Inches(3.5), CW, Inches(0.9), [("기획서 구체화 — Product Deep Dive", 26, CREAM, True)])
text(s, MX, Inches(4.25), CW, Inches(0.6),
     [("\"오늘 점심 뭐 먹지?\"를 3분 안에 끝내는 AI 실시간 그룹 결정 서비스", 15, CREAM_DIM, False)])
text(s, MX, Inches(5.0), CW, Inches(0.5),
     [("페르소나 · 아하 모먼트 · 사용 Flow · 저니맵 · KPI · 설문조사", 13, TERRA, True)])
text(s, MX, SH - Inches(0.9), CW, Inches(0.4),
     [[("스위프 웹 14기 4조  ·  MVP 배포: moddi-demo.vercel.app", 11, CREAM_DIM, False)]])

# ============================================================ 2. AGENDA
s = base("이 문서가 다루는 6가지", kicker="AGENDA")
items = [
    ("01", "페르소나", "누구의 어떤 문제를 푸는가 (JTBD)"),
    ("02", "아하 모먼트", "가치를 처음 체감하는 그 순간 두 개"),
    ("03", "사용 Flow", "상태기계·분기·이탈 방어 (배포 MVP 기준)"),
    ("04", "저니맵", "페르소나별 감정선·페인·기회"),
    ("05", "KPI", "North Star · HEART · AARRR · 이벤트 설계"),
    ("06", "설문조사", "가설 검증 — 컨셉(사전) + 만족(사후)"),
]
y = body_top()
colw = (CW - Inches(0.4)) / 2
for i, (n, t_, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = MX + (colw + Inches(0.4)) * col
    yy = y + (Inches(1.5)) * row
    c = rect(s, x, yy, colw, Inches(1.3), fill=CREAM_SOFT, line_color=LINE, line_w=1)
    text(s, x + Inches(0.25), yy + Inches(0.2), Inches(1.0), Inches(0.6), [(n, 30, TERRA, True)])
    text(s, x + Inches(1.25), yy + Inches(0.22), colw - Inches(1.4), Inches(0.5), [(t_, 19, ESPRESSO, True)])
    text(s, x + Inches(1.25), yy + Inches(0.72), colw - Inches(1.4), Inches(0.5), [(d, 12.5, MUTED, False)])

# ============================================================ 3. PRODUCT
s = base("모디가 푸는 문제, 그리고 지금", kicker="CONTEXT")
y = body_top()
text(s, MX, y, CW, Inches(0.9),
     [[("2~8인이 점심을 정할 때 ", 18, ESPRESSO, False),
       ("결정보다 '후보 모으기·의견 합치기'", 18, TERRA_DK, True),
       ("에서 시간이 샌다.", 18, ESPRESSO, False)]])
# two problem pillars
pw = (CW - Inches(0.4)) / 2
py = y + Inches(1.0)
for i, (h_, d) in enumerate([
    ("콜드 스타트", "빈 화면에서 후보를 사람이 채우는 노동"),
    ("진입 마찰", "가입·로그인 벽에서 참여자 이탈"),
]):
    x = MX + (pw + Inches(0.4)) * i
    rect(s, x, py, pw, Inches(1.15), fill=CREAM_SOFT, line_color=LINE)
    text(s, x + Inches(0.3), py + Inches(0.2), pw - Inches(0.6), Inches(0.5), [(h_, 17, TERRA_DK, True)])
    text(s, x + Inches(0.3), py + Inches(0.68), pw - Inches(0.6), Inches(0.4), [(d, 13, MUTED, False)])
# flow chips
fy = py + Inches(1.5)
text(s, MX, fy, CW, Inches(0.4), [("모디의 해법 — 한 줄 흐름", 13, MUTED, True)])
steps = ["방 생성", "AI 후보 5개", "링크·가입X 입장", "실시간 투표", "자동 결정 🎉"]
cx = MX
chipw = Inches(2.15)
for i, st in enumerate(steps):
    chip(s, cx, fy + Inches(0.45), chipw, st, fill=TERRA if i in (1, 4) else ESPRESSO, size=13, h=Inches(0.5))
    cx += chipw
    if i < len(steps) - 1:
        text(s, cx, fy + Inches(0.45), Inches(0.32), Inches(0.5), [("→", 18, TERRA, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(0.32)
text(s, MX, fy + Inches(1.25), CW, Inches(0.5),
     [[("지금 상태: ", 13, ESPRESSO, True),
       ("동작하는 MVP 배포 완료(강남역 점심 데모) — 이 기획서는 이를 실제 제품으로 구체화한다.", 13, MUTED, False)]])

# ============================================================ 4. PERSONA OVERVIEW
s = base("우리는 누구를 위해 만드는가", kicker="01 · 페르소나")
y = body_top()
cards = [
    ("🅰 김총대 (28)", "방장 · 결정 주도자", "P0 · 1순위", "\"내가 또 정해? 다 같이 30초만 눌러줘\"", "후보 찾는 노동 0"),
    ("🅱 이아무 (24)", "참여자 · 묻어가는 사람", "P0 · 1순위", "\"가입 말고, 링크 주면 바로 누를게\"", "마찰 없는 의견 표출"),
    ("🅲 박모임 (31)", "정기 모임 운영자", "P1 · 확장", "\"매주 똑같은 고민 그만하고 싶어\"", "기록 기반 재추천(리텐션)"),
]
cw3 = (CW - Inches(0.6)) / 3
for i, (name, role, pri, quote, val) in enumerate(cards):
    x = MX + (cw3 + Inches(0.3)) * i
    dim = i == 2
    rect(s, x, y, cw3, Inches(4.2), fill=CREAM_SOFT, line_color=TERRA if not dim else LINE, line_w=1.5 if not dim else 1)
    text(s, x + Inches(0.28), y + Inches(0.28), cw3 - Inches(0.5), Inches(0.5), [(name, 18, ESPRESSO, True)])
    text(s, x + Inches(0.28), y + Inches(0.78), cw3 - Inches(0.5), Inches(0.4), [(role, 13, MUTED, True)])
    chip(s, x + Inches(0.28), y + Inches(1.25), Inches(1.7), pri, fill=TERRA if not dim else MUTED, size=11)
    text(s, x + Inches(0.28), y + Inches(1.95), cw3 - Inches(0.5), Inches(1.2),
         [(quote, 14, TERRA_DK, True)], line_spacing=1.15)
    text(s, x + Inches(0.28), y + Inches(3.35), cw3 - Inches(0.5), Inches(0.4), [("핵심 가치", 11, MUTED, True)])
    text(s, x + Inches(0.28), y + Inches(3.68), cw3 - Inches(0.5), Inches(0.5), [(val, 13.5, ESPRESSO, True)])
text(s, MX, y + Inches(4.35), CW, Inches(0.4),
     [[("설계 우선순위: ", 12, TERRA_DK, True),
       ("P0 MVP는 🅰가 방 만들고 🅱들이 우르르 투표하는 '일회성 빠른 결정'에 집중. 🅲(개인화·리텐션)는 P1.", 12, MUTED, False)]])

# ============================================================ 5. PERSONA DEEP A & B
for (name, sub, jtbd, success, pain, quote) in [
    ("🅰 김총대", "28 · 판교 스타트업 백엔드 3년차 · 팀 점심 주 5회 총대",
     ["기능: 모인 사람이 3분 안에 한 곳에 합의", "감정: '또 나만 정한다'는 부담에서 해방", "사회: '빠르게 정리해주는 사람'으로 보이기"],
     "5분 안에 '안 가본 집'이 모두 동의로 정해짐",
     "후보 수집이 매번 본인 노동 · '아무거나'에 교착 · 못 먹는 사람 뒤늦게 등장",
     "\"내가 또 정해? 그냥 다 같이 30초만 눌러줘.\""),
    ("🅱 이아무", "24 · 경영 3학년 · 모임 식사 주 3~4회 · 설치/가입 강한 거부감",
     ["기능: 설치·가입 없이 선호(특히 비선호) 한 번에", "감정: 까다롭게 안 보이며 '이건 빼줘'를 안전하게", "사회: 민폐 없이 참여하고 한 표 보탠 소속감"],
     "가입 없이 10초 만에 입장해 '이건 빼줘(👎)'를 표시",
     "가입·로그인 벽 · 긴 단톡 스크롤 · 내 '못 먹는 것'이 무시됨",
     "\"가입 말고, 링크 주면 바로 누를게.\""),
]:
    s = base(f"{name} — 딥다이브", kicker="01 · 페르소나")
    y = body_top()
    text(s, MX, y, CW, Inches(0.4), [(sub, 13, MUTED, True)])
    # JTBD card
    jy = y + Inches(0.55)
    rect(s, MX, jy, CW, Inches(1.85), fill=CREAM_SOFT, line_color=LINE)
    text(s, MX + Inches(0.3), jy + Inches(0.18), CW - Inches(0.6), Inches(0.4), [("JTBD — 해결하려는 진짜 일", 13, TERRA_DK, True)])
    text(s, MX + Inches(0.3), jy + Inches(0.62), CW - Inches(0.6), Inches(1.1),
         [[("· ", 14, TERRA, True), (j, 14, ESPRESSO, False)] for j in jtbd], line_spacing=1.25, space_after=4)
    # success + pain row
    ry = jy + Inches(2.1)
    half = (CW - Inches(0.4)) / 2
    rect(s, MX, ry, half, Inches(1.25), fill=None, line_color=TERRA, line_w=1.5)
    text(s, MX + Inches(0.25), ry + Inches(0.18), half - Inches(0.5), Inches(0.4), [("성공의 정의 ✅", 12.5, TERRA_DK, True)])
    text(s, MX + Inches(0.25), ry + Inches(0.6), half - Inches(0.5), Inches(0.6), [(success, 14, ESPRESSO, True)], line_spacing=1.15)
    rect(s, MX + half + Inches(0.4), ry, half, Inches(1.25), fill=None, line_color=LINE, line_w=1.5)
    text(s, MX + half + Inches(0.65), ry + Inches(0.18), half - Inches(0.5), Inches(0.4), [("페인포인트 ⚠️", 12.5, MUTED, True)])
    text(s, MX + half + Inches(0.65), ry + Inches(0.6), half - Inches(0.5), Inches(0.6), [(pain, 13, MUTED, False)], line_spacing=1.15)
    # quote
    text(s, MX, ry + Inches(1.5), CW, Inches(0.6), [(quote, 19, TERRA_DK, True)], align=PP_ALIGN.CENTER)

# ============================================================ 7. AHA
s = base("아하 모먼트 — 가치를 체감하는 그 순간", kicker="02 · 아하 모먼트")
y = body_top()
text(s, MX, y, CW, Inches(0.4), [("역할이 둘이라 아하도 둘. 둘 다 못 터지면 그룹 결정이 시작조차 안 된다.", 14, MUTED, False)])
half = (CW - Inches(0.4)) / 2
ay = y + Inches(0.6)
ahas = [
    ("아하 ①", "방장", "조건만 넣었는데\n후보 5개가 즉시 뜬다", "\"내가 검색 안 해도 되네\"", "방 생성 → 첫 후보", "TTV ≤ 3초", "도달률 ≥ 98%"),
    ("아하 ②", "참여자", "가입창 없이 바로 투표,\n표가 실시간으로 차오른다", "\"가입도 안 했는데 벌써 같이 정하네\"", "링크 클릭 → 첫 투표", "TTV ≤ 10초", "도달률 ≥ 60%"),
]
for i, (tag, who, what, quote, funnel, ttv, rate) in enumerate(ahas):
    x = MX + (half + Inches(0.4)) * i
    rect(s, x, ay, half, Inches(3.95), fill=CREAM_SOFT, line_color=TERRA, line_w=1.5)
    chip(s, x + Inches(0.3), ay + Inches(0.28), Inches(1.2), tag, fill=TERRA, size=13)
    text(s, x + Inches(1.65), ay + Inches(0.3), half - Inches(1.8), Inches(0.4), [(who, 14, MUTED, True)])
    text(s, x + Inches(0.3), ay + Inches(0.95), half - Inches(0.6), Inches(1.0), [(l, 18, ESPRESSO, True) for l in [what]], line_spacing=1.2)
    text(s, x + Inches(0.3), ay + Inches(2.05), half - Inches(0.6), Inches(0.7), [(quote, 15, TERRA_DK, True)], line_spacing=1.15)
    # metrics row
    my = ay + Inches(2.95)
    text(s, x + Inches(0.3), my, half - Inches(0.6), Inches(0.3), [(funnel, 12, MUTED, True)])
    chip(s, x + Inches(0.3), my + Inches(0.4), Inches(1.6), ttv, fill=ESPRESSO, size=12)
    chip(s, x + Inches(2.0), my + Inches(0.4), Inches(1.9), rate, fill=ESPRESSO, size=12)
text(s, MX, ay + Inches(4.05), CW, Inches(0.4),
     [[("철칙: ", 12, TERRA_DK, True), ("아하 전에는 어떤 벽도 없다(가입·로그인 금지) · 빈 후보·표 유실은 아하를 즉사시킨다.", 12, MUTED, False)]])

# ============================================================ 8. FLOW
s = base("사용 Flow — 상태기계 한 장", kicker="03 · 사용 Flow")
y = body_top()
states = [
    ("waiting", "대기", "후보 카드 · [투표 시작] · URL 복사\n참여자는 미리보기(비활성)·'N명 모이는 중'"),
    ("voting", "투표", "후보 👍/👎 · 타이머 · 1초 폴링\n방장만 [지금 마감]"),
    ("decided", "결정", "🎉 1등 카드 + 컨페티\n전원 동시 확인 · 잠김"),
]
bw = (CW - Inches(1.0)) / 3
for i, (en, kr, desc) in enumerate(states):
    x = MX + (bw + Inches(0.5)) * i
    rect(s, x, y + Inches(0.2), bw, Inches(2.0), fill=ESPRESSO if i == 2 else CREAM_SOFT, line_color=TERRA if i == 2 else LINE, line_w=1.5)
    text(s, x + Inches(0.25), y + Inches(0.4), bw - Inches(0.5), Inches(0.45), [[(kr, 18, CREAM if i==2 else ESPRESSO, True), ("  "+en, 11, (CREAM_DIM if i==2 else MUTED), False)]])
    text(s, x + Inches(0.25), y + Inches(0.95), bw - Inches(0.5), Inches(1.1), [(desc, 12, CREAM_DIM if i==2 else MUTED, False)], line_spacing=1.2)
    if i < 2:
        text(s, x + bw, y + Inches(0.85), Inches(0.5), Inches(0.6), [("→", 22, TERRA, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# transitions
text(s, MX, y + Inches(2.45), CW, Inches(0.4),
     [[("전이: ", 12, TERRA_DK, True), ("[방장 투표 시작]  →  [타이머 0 · 자동 우승 공개]  또는  [방장 '지금 마감']  (서버 크론 없이 폴링 lazy 확정)", 12.5, MUTED, False)]])
# key decisions
dy = y + Inches(2.88)
text(s, MX, dy, CW, Inches(0.4), [("핵심 설계 분기", 13, ESPRESSO, True)])
decs = [
    "방장이 시작 눌러야 타이머 ON — 링크 먼저 뿌리고 모인 뒤 시작",
    "동점 → 그 점수에 먼저 도달한 후보 승(결정론적, 수동 개입 0)",
    "방장도 투표자 — 화면 = 참여자 + 운영 버튼",
    "후보 다 별로면 누구나 [+ 후보 직접 추가]",
]
for i, d in enumerate(decs):
    yy = dy + Inches(0.4) + Inches(0.42) * i
    text(s, MX, yy, CW, Inches(0.42), [[("· ", 13, TERRA, True), (d, 13, MUTED, False)]])

# ============================================================ 9. JOURNEY
s = base("저니맵 — 감정선과 아하 정점", kicker="04 · 저니맵")
y = body_top()
def journey_row(s, label, stages, ytop, peak_idx):
    text(s, MX, ytop, Inches(2.2), Inches(0.4), [(label, 14, ESPRESSO, True)])
    n = len(stages); gw = CW / n
    for i, (emo, st) in enumerate(stages):
        x = MX + gw * i
        is_peak = i == peak_idx
        rect(s, x + Inches(0.05), ytop + Inches(0.45), gw - Inches(0.1), Inches(1.15),
             fill=TERRA if is_peak else CREAM_SOFT, line_color=TERRA if is_peak else LINE, line_w=1)
        text(s, x + Inches(0.05), ytop + Inches(0.58), gw - Inches(0.1), Inches(0.5), [(emo, 22, WHITE if is_peak else ESPRESSO, False)], align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.05), ytop + Inches(1.12), gw - Inches(0.1), Inches(0.45), [(st, 10.5, CREAM if is_peak else MUTED, is_peak)], align=PP_ALIGN.CENTER, line_spacing=1.0)
journey_row(s, "🅰 김총대", [("😩","트리거"),("🤔","폼 입력"),("😀","후보 즉시"),("🙂","초대"),("😌","투표 관찰"),("😄","결정")], y, 2)
text(s, MX + CW*2/6, y + Inches(0.3), CW/6, Inches(0.3), [("★ 아하①", 10, TERRA_DK, True)], align=PP_ALIGN.CENTER)
y2 = y + Inches(1.9)
journey_row(s, "🅱 이아무", [("😐","링크 받음"),("🙂","클릭·안도"),("😀","가입X 입장"),("🙂","후보 훑기"),("😄","투표"),("😊","결과")], y2, 4)
text(s, MX + CW*4/6, y2 + Inches(0.3), CW/6, Inches(0.3), [("★ 아하②", 10, TERRA_DK, True)], align=PP_ALIGN.CENTER)
text(s, MX, y2 + Inches(1.85), CW, Inches(0.4),
     [[("관찰: ", 12, TERRA_DK, True), ("아하①·② 모두 감정 곡선의 가장 가파른 상승 구간 — 흐려지면 전체 경험이 평탄해진다.", 12, MUTED, False)]])

# ============================================================ 10. KPI
s = base("KPI — 무엇을 측정하는가", kicker="05 · KPI")
y = body_top()
# North star band
rect(s, MX, y, CW, Inches(1.15), fill=ESPRESSO)
text(s, MX + Inches(0.35), y + Inches(0.2), CW - Inches(0.7), Inches(0.4), [("North Star Metric", 13, TERRA, True)])
text(s, MX + Inches(0.35), y + Inches(0.55), CW - Inches(0.7), Inches(0.5),
     [[("주간 '완료된 그룹 결정' 수 ", 20, CREAM, True), ("(decided & 실투표자 ≥ 2명)", 14, CREAM_DIM, False)]])
# HEART row
hy = y + Inches(1.4)
text(s, MX, hy, CW, Inches(0.35), [("HEART — 경험 품질", 13, ESPRESSO, True)])
heart = [("Happiness","만족 ≥4.0"),("Engagement","방당 투표자 ≥3"),("Adoption","아하 98%/60%"),("Retention","재사용 의향"),("Task success","완료율 ≥80%")]
hw = (CW - Inches(0.6)) / 5
for i,(k,v) in enumerate(heart):
    x = MX + (hw + Inches(0.15))*i
    rect(s, x, hy+Inches(0.4), hw, Inches(1.05), fill=CREAM_SOFT, line_color=LINE)
    text(s, x+Inches(0.12), hy+Inches(0.55), hw-Inches(0.24), Inches(0.5), [(k, 12.5, TERRA_DK, True)], align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.12), hy+Inches(1.0), hw-Inches(0.24), Inches(0.4), [(v, 11.5, ESPRESSO, True)], align=PP_ALIGN.CENTER)
# funnel
fy = hy + Inches(1.75)
text(s, MX, fy, CW, Inches(0.35), [("AARRR 활성화 퍼널", 13, ESPRESSO, True)])
fun = ["링크 오픈","입장","첫 투표(아하②)","결정 완료"]
fw = (CW - Inches(1.0))/4
for i,f in enumerate(fun):
    x = MX + (fw+Inches(0.33))*i
    chip(s, x, fy+Inches(0.42), fw, f, fill=TERRA if i in (2,3) else ESPRESSO, size=12, h=Inches(0.5))
    if i<3:
        text(s, x+fw, fy+Inches(0.42), Inches(0.33), Inches(0.5), [("▸", 16, TERRA, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================ 11. SCOREBOARD
s = base("데모데이 스코어보드 — 발표 때 띄울 숫자", kicker="05 · KPI")
y = body_top()
sb = [("결정 완료율","방 생성 → 확정 도달","≥ 80%"),
      ("결정 소요시간","방 생성 → 확정","≤ 3분"),
      ("방당 참여자","투표 참여 인원","≥ 3명"),
      ("AI 후보 채택률","최종 선택이 AI 후보","측정·증명")]
cw2 = (CW - Inches(0.6))/2
for i,(k,d,v) in enumerate(sb):
    col=i%2; row=i//2
    x = MX + (cw2+Inches(0.6))*col
    yy = y + (Inches(1.7))*row
    rect(s, x, yy, cw2, Inches(1.45), fill=CREAM_SOFT, line_color=TERRA if i in (0,1) else LINE, line_w=1.5 if i in (0,1) else 1)
    text(s, x+Inches(0.3), yy+Inches(0.22), cw2-Inches(2.0), Inches(0.5), [(k, 17, ESPRESSO, True)])
    text(s, x+Inches(0.3), yy+Inches(0.72), cw2-Inches(2.0), Inches(0.5), [(d, 12, MUTED, False)])
    text(s, x+cw2-Inches(1.9), yy+Inches(0.35), Inches(1.7), Inches(0.7), [(v, 22, TERRA_DK, True)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
text(s, MX, y+Inches(3.5), CW, Inches(0.4),
     [[("핵심: ", 12, TERRA_DK, True), ("정밀 대시보드보다 '심사위원이 직접 만든 방의 숫자'가 설득력 — 현장 참여형 데모와 결합.", 12, MUTED, False)]])

# ============================================================ 12. SURVEY
s = base("설문조사 — 가설을 데이터로", kicker="06 · 설문조사")
y = body_top()
# hypotheses
text(s, MX, y, CW, Inches(0.35), [("검증 가설", 13, ESPRESSO, True)])
hyps = [("H1","첫 방에서 '결정 완료' 경험 → 재사용 의향 ↑"),
        ("H2","방당 실투표 ≥ 3명 → 만족·재사용 ↑"),
        ("H3","그룹 점심 결정은 충분히 자주·고통스러운 문제다")]
for i,(h,t_) in enumerate(hyps):
    yy = y+Inches(0.42)+Inches(0.5)*i
    chip(s, MX, yy, Inches(0.7), h, fill=TERRA, size=12, h=Inches(0.36))
    text(s, MX+Inches(0.9), yy, CW-Inches(0.9), Inches(0.4), [(t_, 13.5, ESPRESSO, False)], anchor=MSO_ANCHOR.MIDDLE)
# two surveys
sy = y + Inches(2.3)
half = (CW - Inches(0.4))/2
for i,(tag, who, items_) in enumerate([
    ("설문 A · 컨셉 검증", "사전 · 잠재 사용자 20~30명", ["문제 빈도·심각도(H3)","현재 대안 & 불편","가입 거부감·프라이버시","컨셉 사용 의향"]),
    ("설문 B · 만족·재사용", "사후 · 사용 직후 5~10명", ["전반 만족도(5점)","재사용 의향(NPS-lite)","아하 체감 2종","개방형: 좋았던/불편한 점"]),
]):
    x = MX + (half+Inches(0.4))*i
    rect(s, x, sy, half, Inches(2.3), fill=CREAM_SOFT, line_color=LINE)
    text(s, x+Inches(0.28), sy+Inches(0.2), half-Inches(0.56), Inches(0.4), [(tag, 15, TERRA_DK, True)])
    text(s, x+Inches(0.28), sy+Inches(0.62), half-Inches(0.56), Inches(0.35), [(who, 11.5, MUTED, True)])
    text(s, x+Inches(0.28), sy+Inches(1.0), half-Inches(0.56), Inches(1.2),
         [[("· ", 12, TERRA, True), (it, 12.5, ESPRESSO, False)] for it in items_], line_spacing=1.3, space_after=3)

# ============================================================ 13. CLOSING
s = prs.slides.add_slide(BLANK); bg(s, ESPRESSO)
rect(s, Inches(0), SH-Inches(0.32), SW, Inches(0.32), fill=TERRA, radius=0)
text(s, MX, Inches(1.0), CW, Inches(0.4), [("NEXT", 14, TERRA, True)])
text(s, MX, Inches(1.4), CW, Inches(0.8), [("기획서 → 실행으로", 32, CREAM, True)])
road = [
    ("P0 (지금)", "강남역 점심 1카테고리 · 게스트 입장 · 실시간 투표 · 자동 결정 — 배포 완료"),
    ("P0 검증", "데모데이 스코어보드 + 설문 A/B로 H1·H2·H3 검증"),
    ("P1", "개인화 추천(🅲 리텐션 루프) · 지도 연동 · AI 타이브레이크"),
]
ry = Inches(2.6)
for i,(p,d) in enumerate(road):
    yy = ry + Inches(1.0)*i
    chip(s, MX, yy, Inches(1.5), p, fill=TERRA, size=13, h=Inches(0.42))
    text(s, MX+Inches(1.75), yy, CW-Inches(1.75), Inches(0.7), [(d, 14.5, CREAM_DIM, False)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
text(s, MX, SH-Inches(1.15), CW, Inches(0.5),
     [[("모디", 18, CREAM, True), (".", 18, TERRA, True), ("  같이 정하자  ·  moddi-demo.vercel.app", 13, CREAM_DIM, False)]])

out = "docs/planning/모디_기획서_구체화.pptx"
prs.save(out)
print(f"saved {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
